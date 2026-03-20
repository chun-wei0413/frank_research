#!/usr/bin/env python3
"""Rebuild skills.pptx - condensed 30-min version with all key PDF content."""

from pptx import Presentation
from pptx.util import Pt

INPUT = 'doc/skills.pptx'
OUTPUT = 'doc/skills.pptx'

prs = Presentation(INPUT)

LAYOUT_TITLE = 0
LAYOUT_OBJECT = 1
LAYOUT_SECTION = 2
LAYOUT_TWO_OBJ = 3

# Remove all slides except the first one (cover)
total = len(prs.slides)
for idx in range(total - 1, 0, -1):
    rId = prs.slides._sldIdLst[idx].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
    )
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[idx])

print(f"After cleanup: {len(prs.slides)} (cover only)")


def add_object_slide(title_text, bullets, sub_bullets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_OBJECT])
    slide.placeholders[0].text = title_text
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = sb
                sp.level = 1
                sp.font.size = Pt(14)
    return slide


def add_section_slide(title_text, subtitle_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_SECTION])
    slide.placeholders[0].text = title_text
    if subtitle_text:
        slide.placeholders[1].text = subtitle_text
    return slide


def add_two_col_slide(title_text, left_bullets, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TWO_OBJ])
    slide.placeholders[0].text = title_text
    tf_left = slide.placeholders[1].text_frame
    tf_left.clear()
    for i, b in enumerate(left_bullets):
        p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(14)
    tf_right = slide.placeholders[2].text_frame
    tf_right.clear()
    for i, b in enumerate(right_bullets):
        p = tf_right.paragraphs[0] if i == 0 else tf_right.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(14)
    return slide


# ============================================================
# Slide 2: Outline (~1 min)
# ============================================================
add_object_slide("Outline", [
    "1. 什麼是 Skill？ — 定義、結構與核心設計原則",
    "2. Skills + MCP — 知識層與連接層的協作關係",
    "3. 規劃與設計 — 使用案例、YAML Frontmatter、指令撰寫",
    "4. 測試與迭代 — 成功標準、測試方法、Skill-Creator",
    "5. 發布與設計模式 — 發布方式、五大模式、問題排除",
])

# ============================================================
# Slide 3: 什麼是 Skill？ (~2 min)
# ============================================================
add_object_slide("什麼是 Skill？", [
    "一組打包成資料夾的指令集，教導 Claude 處理特定任務",
    "教 Claude 一次，之後每次對話都能受益",
    "適用場景：可重複的工作流程",
    "本指南的兩條路徑：",
], sub_bullets={
    2: [
        "從規格生成前端設計 / 以一致方法進行研究",
        "建立遵循團隊風格指南的文件 / 協調多步驟流程",
    ],
    3: [
        "獨立 Skill 開發者 → 聚焦 Fundamentals + Planning (Cat. 1-2)",
        "MCP 整合開發者 → 聚焦 Skills + MCP 段落 (Cat. 3)",
    ]
})

# ============================================================
# Slide 4: 核心設計原則 — 漸進式揭露 + 可組合/可攜 (~2 min)
# ============================================================
add_object_slide("核心設計原則", [
    "漸進式揭露（Progressive Disclosure）：",
    "可組合性：可同時載入多個 Skill 協同運作",
    "可攜性：Claude.ai、Claude Code、API 上運作完全一致",
], sub_bullets={
    0: [
        "第一層 YAML Frontmatter — 永遠載入，判斷何時使用",
        "第二層 SKILL.md 主體 — 相關時載入，完整指令",
        "第三層 references/ 連結檔 — 按需載入，深度參考",
        "效益：最小化 Token 使用量，同時保持專業能力",
    ],
})

# ============================================================
# Slide 5: Skills + MCP 廚房比喻 (~3 min)
# ============================================================
add_two_col_slide("Skills + MCP：廚房比喻", [
    "MCP = 專業廚房（連接層）",
    "• 連接 Claude 到你的服務",
    "• 提供即時資料存取和工具調用",
    "• 回答：Claude 能做什麼？",
    "",
    "沒有 Skills 時：",
    "• 使用者不知道下一步",
    "• 每次對話從零開始",
    "• 結果不一致",
    "• 使用者歸咎於連接器",
], [
    "Skills = 食譜（知識層）",
    "• 教導 Claude 如何使用服務",
    "• 捕捉工作流程和最佳實踐",
    "• 回答：Claude 該怎麼做？",
    "",
    "有 Skills 時：",
    "• 預建工作流程自動啟動",
    "• 一致且可靠的工具使用",
    "• 最佳實踐嵌入每次互動",
    "• 降低整合的學習曲線",
])

# ============================================================
# Slide 6: 三大使用案例類別 + 範例 (~3 min)
# ============================================================
add_object_slide("三大使用案例類別", [
    "Category 1：文件與資產建立",
    "Category 2：工作流程自動化",
    "Category 3：MCP 增強",
    "",
    "建構使用案例範本（以 Sprint Planning 為例）：",
], sub_bullets={
    0: [
        "一致輸出（文件、簡報、程式碼）/ 技巧：風格指南、範本、品質檢查清單",
    ],
    1: [
        "多步驟一致方法論（如 skill-creator）/ 技巧：驗證關卡、迭代精煉",
    ],
    2: [
        "為 MCP 工具加上引導（如 sentry-code-review）/ 技巧：多 MCP 協調、領域知識",
    ],
    4: [
        "Trigger:「help me plan this sprint」",
        "Steps: 取得專案狀態 → 分析產能 → 建議優先順序 → 建立任務",
        "Result: 完整的 Sprint 計畫與已建立的任務",
    ]
})

# ============================================================
# Slide 7: 檔案結構 + YAML (~2 min)
# ============================================================
add_two_col_slide("Skill 的檔案結構與 YAML Frontmatter", [
    "資料夾結構：",
    "  your-skill-name/",
    "  ├── SKILL.md  (必要)",
    "  ├── scripts/    (選用)",
    "  ├── references/ (選用)",
    "  └── assets/     (選用)",
    "",
    "關鍵規則：",
    "• SKILL.md 大小寫敏感",
    "• 資料夾用 kebab-case",
    "• 不放 README.md",
], [
    "最小 YAML Frontmatter：",
    "  ---",
    "  name: your-skill-name",
    "  description: 做什麼 + 何時用",
    "  ---",
    "",
    "選用欄位：license, compatibility,",
    "  allowed-tools, metadata",
    "",
    "安全限制：",
    "• 禁止 XML 標籤 (< >)",
    '• 名稱不可含 "claude"/"anthropic"',
])

# ============================================================
# Slide 8: Description 好壞範例 (~2 min)
# ============================================================
add_two_col_slide("Description：Skill 成敗的關鍵", [
    "好的 Description：",
    "",
    "• 具體說明「做什麼」",
    "• 包含「何時使用」的觸發詞",
    "• 提及相關檔案類型",
    "",
    "範例：",
    '"Manages Linear workflows',
    "including sprint planning.",
    "Use when user mentions",
    '  sprint or Linear tasks."',
], [
    "壞的 Description：",
    "",
    '• "Helps with projects"',
    "   → 太模糊",
    '• "Creates documentation systems"',
    "   → 沒有觸發條件",
    '• "Implements entity model"',
    "   → 太技術化",
    "",
    "原則：使用者怎麼說，",
    "description 就怎麼寫",
])

# ============================================================
# Slide 9: 撰寫指令的最佳實踐 (~2 min)
# ============================================================
add_object_slide("撰寫 SKILL.md 指令的最佳實踐", [
    "具體且可操作",
    "善用漸進式揭露",
    "包含錯誤處理與範例",
    "進階技巧：用腳本做確定性驗證，而非語言指令",
    "對抗模型偷懶：加入 Performance Notes 鼓勵語",
], sub_bullets={
    0: [
        '好："Run python scripts/validate.py --input {filename}"',
        '壞："Validate the data before proceeding"',
    ],
    1: ["核心指令放 SKILL.md、詳細文件移 references/"],
    2: ["MCP 連線失敗步驟、API key 驗證、常見錯誤碼"],
    4: [
        '"Take your time to do this thoroughly"',
        "注意：加在使用者提示中比寫在 SKILL.md 更有效",
    ],
})

# ============================================================
# Slide 10: 定義成功標準 (~2 min)
# ============================================================
add_two_col_slide("定義成功標準", [
    "量化指標：",
    "• Skill 在 90% 相關查詢觸發",
    "  → 執行 10-20 個測試查詢驗證",
    "• 在 X 次工具呼叫內完成流程",
    "  → 與未使用 Skill 時比較",
    "• 每個流程 0 次失敗的 API 呼叫",
    "  → 監控 MCP 伺服器日誌",
], [
    "質化指標：",
    "• 使用者不需提示下一步",
    "• 工作流程無需修正即完成",
    "• 各工作階段結果一致",
    "• 新使用者能否首次就完成？",
    "",
    "效能對比範例：",
    "Without: 15 來回 / 3 失敗 / 12K tokens",
    "With: 自動執行 / 0 失敗 / 6K tokens",
])

# ============================================================
# Slide 11: 測試的三個面向 (~2 min)
# ============================================================
add_object_slide("測試的三個面向", [
    "1. 觸發測試 — Skill 是否在正確時機載入？",
    "2. 功能測試 — 輸出正確？API 成功？錯誤處理正常？",
    "3. 效能比較 — 對比基準（訊息來回、失敗呼叫、Token 消耗）",
    "",
    "測試方法：手動(Claude.ai) / 腳本(Claude Code) / 程式化(Skills API)",
    "Pro Tip：先針對單一困難任務迭代至成功，再萃取為 Skill",
], sub_bullets={
    0: [
        "✅ 明確任務觸發 / ✅ 改述請求觸發 / ❌ 無關主題不觸發",
    ],
})

# ============================================================
# Slide 12: Skill-Creator 與迭代 (~2 min)
# ============================================================
add_two_col_slide("Skill-Creator 與迭代改進", [
    "Skill-Creator 工具：",
    "• 從自然語言描述生成 Skill",
    "• 產生正確格式的 SKILL.md",
    "• 審查並建議改善",
    "",
    "使用方式：",
    '"Use skill-creator to help',
    '  me build a skill for [用途]"',
    "",
    "注意：協助設計與精煉，",
    "不執行自動化測試套件",
], [
    "迭代信號與對策：",
    "",
    "觸發不足：Skill 該載入卻沒載入",
    "→ 加更多觸發詞/技術用語",
    "",
    "觸發過度：無關查詢也載入",
    "→ 加負面觸發 / 縮小範圍",
    "",
    "執行不一致：結果不穩定",
    "→ 改善指令 / 加錯誤處理",
    "",
    "Skills 是活文件，持續迭代",
])

# ============================================================
# Slide 13: 發布與分享 (~2 min)
# ============================================================
add_object_slide("發布與分享", [
    "個人：下載資料夾 → 壓縮 → 上傳 Claude.ai 或放入 Claude Code 目錄",
    "組織：管理員在工作區部署，自動更新、集中管理",
    "API：/v1/skills 端點、Messages API 整合、Agent SDK 支援",
    "開放標準：如同 MCP，Skill 可跨工具和平台使用",
    "",
    "定位原則：聚焦成果而非功能",
], sub_bullets={
    5: [
        '好：「幾秒內建立完整工作區，取代 30 分鐘手動設定」',
        '壞：「包含 YAML frontmatter 和 Markdown 指令的資料夾」',
    ]
})

# ============================================================
# Slide 14: Problem-first vs. Tool-first + 五大模式 (~3 min)
# ============================================================
add_object_slide("五大設計模式", [
    "選擇方法：Problem-first vs. Tool-first（Home Depot 比喻）",
    "",
    "Pattern 1：循序式工作流程 — 特定順序的多步驟流程",
    "Pattern 2：多 MCP 協調 — 跨多個服務的工作流程",
    "Pattern 3：迭代精煉 — 品質因迭代而改善",
    "Pattern 4：上下文感知選擇 — 依情境選不同工具",
    "Pattern 5：領域專業智慧 — 嵌入超越工具的專業知識",
], sub_bullets={
    0: [
        "Problem-first：使用者描述目標 → Skill 編排正確的 MCP 呼叫",
        "Tool-first：使用者有 MCP 存取 → Skill 教導最佳工作流程",
    ],
    2: ["範例：客戶上架 → 建帳號 → 設定付款 → 建訂閱 → 發歡迎信"],
    3: ["範例：Figma 匯出 → Drive 儲存 → Linear 建任務 → Slack 通知"],
    7: ["範例：付款處理前先跑合規檢查、制裁名單、稽核紀錄"],
})

# ============================================================
# Slide 15: Troubleshooting (~2 min)
# ============================================================
add_two_col_slide("常見問題排除", [
    "上傳失敗：",
    "• SKILL.md 大小寫敏感",
    "• YAML 需 --- 分隔符號",
    "• 名稱用 kebab-case",
    "",
    "觸發問題：",
    "• 不足 → 加觸發詞到 description",
    "• 過度 → 加負面觸發、縮小範圍",
    '• 除錯：問 Claude "When would',
    '  you use [skill name]?"',
], [
    "指令未遵循：",
    "• 太冗長 → 精簡、條列式",
    "• 關鍵指令被埋 → 放最上方",
    "• 語言模糊 → 用腳本驗證",
    "",
    "效能問題：",
    "• SKILL.md < 5,000 字",
    "• 詳細文件移 references/",
    "• 避免同時啟用 >20 個 Skill",
    "",
    "MCP 問題：先單獨測 MCP",
])

# ============================================================
# Slide 16: Quick Checklist + 資源 (~1 min)
# ============================================================
add_object_slide("Quick Checklist & 資源", [
    "開始前：辨識 2-3 使用案例、規劃資料夾結構",
    "開發中：kebab-case 名稱、SKILL.md、正確 YAML、清楚指令",
    "上傳前：觸發 + 功能測試通過、壓縮為 .zip",
    "上傳後：真實對話測試、監控觸發、收集回饋迭代",
    "",
    "資源：",
], sub_bullets={
    5: [
        "skill-creator：內建工具，快速生成和審查 Skill",
        "github.com/anthropics/skills：公開範例倉庫",
        "支援：Claude Developers Discord / GitHub Issues",
    ]
})

# ============================================================
# Slide 17: 結論 (~2 min)
# ============================================================
add_object_slide("結論", [
    "Skill 將「一次性的提示工程」封裝為「可重複使用的知識資產」",
    "三層漸進式揭露：平衡 Token 效率與專業深度",
    "MCP + Skills = 連接能力 + 工作流程知識 → 完整的 AI 解決方案",
    "Description 是成敗關鍵 — 決定 Skill 何時被觸發",
    "持續迭代：Skill 是活文件，透過測試與回饋不斷改善",
    "",
    "從今天開始：用 skill-creator 花 15-30 分鐘建立你的第一個 Skill",
])

# ============================================================
# Final Slide: Thanks for listening
# ============================================================
thanks = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
thanks.placeholders[0].text = "Thanks for listening"

print(f"Final slides: {len(prs.slides)}")
prs.save(OUTPUT)
print(f"Saved to {OUTPUT}")
